import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_xml():
    f = "sourcefiles/aicc/Introducing ChatGPT and generative AI at Oxford - draft version.pptx"
    prs = Presentation(f)
    
    # Check Slide 27
    slide = prs.slides[26] # 0-indexed
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            print(f"Checking Media on Slide 27...")
            # Dump XML
            print(shape.element.xml)
            
            # Look for r:link or similar
            # For video, it's usually <p:videoFile r:link="rIdX"/> inside <p:nvPicPr><p:nvPr><p:videoFile ...
            # Or <a:videoFile>
            
            # Let's inspect namespaces
            ns = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            }
            
            # The shape is likely a Picture frame with video props
            # Try to find videoFile
            
            # Common path: p:pic -> p:nvPicPr -> p:nvPr -> a:videoFile (or p:videoFile)
            video_file = shape.element.find(".//p:videoFile", ns)
            if video_file is not None:
                rid = video_file.get(f"{{{ns['r']}}}link")
                print(f"Found video rId: {rid}")
                
                if rid:
                    try:
                        part = shape.part.related_part(rid)
                        print(f"Part Content Type: {part.content_type}")
                        print(f"Blob size: {len(part.blob)} bytes")
                    except Exception as e:
                        print(f"Error getting part: {e}")
            else:
                print("No p:videoFile found")

analyze_xml()
