import os
import logging
from datetime import datetime
from typing import List, Dict, Optional
from pptx import Presentation as PptxPresentation
# MSO_SHAPE_TYPE and PP_PLACEHOLDER not needed - we check for image.blob directly
from ..core.models import (
    Presentation, PresentationMetadata, Section, Slide,
    ContentBlock, HeadingBlock, ParagraphBlock, ListBlock, ListItem,
    ImageBlock, SmartArtBlock, SmartArtNode, TableBlock, LinkBlock, VideoBlock,
    ShapeBlock, TextRun
)
from ..core.interfaces import IPresentationExtractor
from .smartart_extractor import SmartArtExtractor

logger = logging.getLogger(__name__)

class PptxExtractor(IPresentationExtractor):
    def __init__(self):
        self.smartart_extractor = SmartArtExtractor()

    def extract(self, file_path: str, media_output_dir: str) -> Presentation:
        filename = os.path.basename(file_path)
        file_id = os.path.splitext(filename)[0]
        
        file_media_dir = os.path.join(media_output_dir, file_id)
        os.makedirs(file_media_dir, exist_ok=True)

        logger.info(f"Parsing {file_path}")
        prs = PptxPresentation(file_path)

        # 1. Sections
        sections_map = self._get_sections_mapping(prs)
        
        # 2. Iterate
        processed_sections = []
        slide_count = 0
        image_count = 0
        
        if not sections_map:
            sections_map = [{"title": "Default", "slides": list(prs.slides)}]

        for sec_data in sections_map:
            section_slides = []
            for pptx_slide in sec_data['slides']:
                slide_count += 1
                slide_obj, imgs = self._process_slide(pptx_slide, slide_count, file_id, file_media_dir)
                section_slides.append(slide_obj)
                image_count += imgs
            
            processed_sections.append(Section(
                title=sec_data['title'],
                slides=section_slides
            ))

        # 3. Metadata
        metadata = PresentationMetadata(
            id=file_id,
            source_file=filename,
            processed_at=datetime.now().isoformat(),
            stats={"slide_count": slide_count, "image_count": image_count}
        )

        return Presentation(metadata=metadata, sections=processed_sections)

    def _get_sections_mapping(self, prs) -> list:
        mapping = []
        # 1. Native Sections
        if hasattr(prs, 'sections') and len(prs.sections) > 0:
            try:
                for section in prs.sections:
                    slides = list(section.slides)
                    if slides:
                        mapping.append({"title": section.title, "slides": slides})
                return mapping
            except Exception as e:
                logger.warning(f"Error accessing native sections: {e}")
        
        # 2. XML Fallback
        try:
            slide_id_to_slide = {slide.slide_id: slide for slide in prs.slides}
            nsmap = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main'
            }
            ext_lst = prs.element.find("./p:extLst", namespaces=nsmap)
            if ext_lst is not None:
                for ext in ext_lst.findall("./p:ext", namespaces=nsmap):
                    section_list_xml = ext.find(".//p14:sectionLst", namespaces=nsmap)
                    if section_list_xml is not None:
                        for section in section_list_xml.findall("./p14:section", namespaces=nsmap):
                            section_name = section.get("name")
                            section_slides = []
                            sld_id_lst = section.find("./p14:sldIdLst", namespaces=nsmap)
                            if sld_id_lst is not None:
                                for sld_id_tag in sld_id_lst.findall("./p14:sldId", namespaces=nsmap):
                                    sid = int(sld_id_tag.get("id"))
                                    if sid in slide_id_to_slide:
                                        section_slides.append(slide_id_to_slide[sid])
                            if section_slides:
                                mapping.append({"title": section_name, "slides": section_slides})
                        return mapping
        except Exception as e:
            logger.warning(f"XML section extraction failed: {e}")

        return []

    def _extract_emf_embedded_image(self, emf_data: bytes) -> tuple | None:
        """
        Extract embedded JPEG/PNG image from EMF+ (Enhanced Metafile Plus) format.

        EMF+ files often contain raster images embedded in GDIC comment records.
        This is common when screenshots are pasted into PowerPoint.

        Returns:
            Tuple of (image_bytes, extension) if found, None otherwise
        """
        import struct

        # EMF+ uses comment records (type 70) to store GDI+ data
        # Look for GDIC records which often contain embedded images
        pos = 0
        while pos < len(emf_data) - 8:
            try:
                record_type, record_size = struct.unpack('<II', emf_data[pos:pos+8])
            except struct.error:
                break

            if record_type == 70:  # EMR_COMMENT (may contain EMF+ or GDIC data)
                comment_data = emf_data[pos+8:pos+record_size]
                if len(comment_data) > 8:
                    # Check for GDIC identifier (contains embedded images)
                    identifier = comment_data[4:8]
                    if identifier == b'GDIC':
                        # Search for JPEG signature (FFD8FF)
                        jpg_sig = b'\xff\xd8\xff'
                        jpg_pos = comment_data.find(jpg_sig)
                        if jpg_pos >= 0:
                            # Extract JPEG - find EOI marker (0xFFD9)
                            jpg_data = comment_data[jpg_pos:]
                            eoi_pos = jpg_data.find(b'\xff\xd9')
                            if eoi_pos > 0:
                                jpg_data = jpg_data[:eoi_pos+2]
                                return (jpg_data, 'jpg')

                        # Search for PNG signature
                        png_sig = b'\x89PNG\r\n\x1a\n'
                        png_pos = comment_data.find(png_sig)
                        if png_pos >= 0:
                            # Extract PNG - find IEND chunk
                            png_data = comment_data[png_pos:]
                            iend_pos = png_data.find(b'IEND')
                            if iend_pos > 0:
                                png_data = png_data[:iend_pos+8]  # Include IEND + CRC
                                return (png_data, 'png')

            if record_type == 14 or record_size == 0:  # EOF or invalid
                break
            pos += record_size

        return None

    def _process_slide(self, slide, order: int, file_id: str, media_dir: str):
        content = []
        notes = ""
        image_count = 0

        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"

        # Extract animation order map (shape_id -> animation_order)
        animation_map = self._extract_animation_map(slide)

        title = ""
        title_url = None
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()
            # Extract hyperlink from title before adding heading
            if slide.shapes.title.has_text_frame:
                title_url = self._extract_hyperlink_from_text_frame(slide.shapes.title.text_frame)
            content.append(HeadingBlock(text=title, level=1))
            # If title has a URL, add it as a link block too
            if title_url:
                content.append(LinkBlock(text=title, url=title_url))

        shapes = list(slide.shapes)
        # Handle None values in top/left - some shapes don't have position
        shapes.sort(key=lambda s: (s.top if (hasattr(s, 'top') and s.top is not None) else 0,
                                   s.left if (hasattr(s, 'left') and s.left is not None) else 0))

        for shape in shapes:
            # Check for embedded video in ANY shape type first
            video_info = self._extract_video_info(shape, order, file_id, media_dir)
            if video_info:
                if 'url' in video_info:
                    # External video (YouTube, etc.)
                    content.append(LinkBlock(text=video_info['title'], url=video_info['url']))
                elif 'src' in video_info:
                    # Embedded video (MP4, etc.)
                    content.append(VideoBlock(src=video_info['src'], title=video_info['title']))
                # Continue to also process other content in this shape if any

            if shape == slide.shapes.title:
                continue

            # Table
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append(cell.text_frame.text.strip())
                    rows.append(cells)
                content.append(TableBlock(rows=rows))
                continue

            # SmartArt
            is_graphic_frame = shape.shape_type == 6 or type(shape).__name__ == 'GraphicFrame'
            if is_graphic_frame:
                sa_data = self.smartart_extractor.extract(shape, slide.part, file_id, media_dir)
                if sa_data:
                    def dict_to_node(d):
                        return SmartArtNode(
                            id=d['id'], text=d['text'], level=d['level'], icon=d['icon'], icon_alt=d.get('icon_alt'),
                            children=[dict_to_node(c) for c in d['children']]
                        )
                    nodes = [dict_to_node(n) for n in sa_data['nodes']]
                    content.append(SmartArtBlock(layout=sa_data['layout'], nodes=nodes))
                    continue

            # Image - extract from ANY shape that has image.blob, not just PICTURE type
            # (Placeholder shapes with images are type 14, not 13, but still have image data)
            if hasattr(shape, "image") and hasattr(shape.image, "blob"):
                try:
                    blob = shape.image.blob
                    # Get the actual image extension from python-pptx
                    ext = getattr(shape.image, 'ext', 'png')
                    content_type = getattr(shape.image, 'content_type', 'image/png')

                    # Handle EMF/WMF vector formats - try to extract embedded images or convert
                    if ext in ('emf', 'wmf') or content_type in ('image/x-emf', 'image/x-wmf'):
                        converted = False
                        # First try to extract embedded images from EMF+ format
                        try:
                            embedded = self._extract_emf_embedded_image(blob)
                            if embedded:
                                blob, ext = embedded
                                converted = True
                                logger.info(f"Extracted embedded image from EMF+ format")
                        except Exception as emf_err:
                            logger.debug(f"EMF+ extraction failed: {emf_err}")

                        # Fallback to Pillow conversion
                        if not converted:
                            try:
                                from PIL import Image
                                import io
                                img = Image.open(io.BytesIO(blob))
                                png_buffer = io.BytesIO()
                                img.save(png_buffer, format='PNG')
                                blob = png_buffer.getvalue()
                                ext = 'png'
                                logger.info(f"Converted {content_type} to PNG")
                            except Exception as conv_err:
                                logger.warning(f"Could not convert {ext} to PNG: {conv_err}")

                    image_filename = f"slide_{order}_{shape.shape_id}.{ext}"
                    image_path = os.path.join(media_dir, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(blob)

                    src = f"media/{file_id}/{image_filename}"
                    alt = shape.name if hasattr(shape, 'name') else "Slide Image"
                    content.append(ImageBlock(src=src, alt=alt))
                    image_count += 1
                    continue
                except Exception as e:
                    logger.debug(f"Could not extract image from shape: {e}")

            # Auto shapes (arrows, connectors, symbols, etc.)
            shape_block = self._extract_auto_shape(shape, slide, animation_map)
            if shape_block:
                content.append(shape_block)
                # If shape has text, also process as text (don't skip)
                if shape.has_text_frame and shape.text_frame.text.strip():
                    blocks = self._process_text_frame(shape.text_frame)
                    content.extend(blocks)
                continue

            # Text
            if shape.has_text_frame:
                if not shape.text_frame.text.strip():
                    continue
                blocks = self._process_text_frame(shape.text_frame)
                content.extend(blocks)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text

        slide_obj = Slide(
            order=order,
            title=title,
            layout=layout_name,
            notes=notes,
            content=content
        )
        
        return slide_obj, image_count

    def _process_text_frame(self, text_frame) -> List[ContentBlock]:
        """Process text frame and extract list items with formatting."""
        blocks = []
        current_list_items = []

        for p in text_frame.paragraphs:
            text = p.text.strip()
            if not text:
                continue

            # Extract formatted runs
            runs = self._extract_formatted_runs(p)
            url = self._extract_hyperlink_from_paragraph(p)
            current_list_items.append(ListItem(text=text, level=p.level, url=url, runs=runs))

        if current_list_items:
            blocks.append(ListBlock(style="bullet", items=current_list_items))

        return blocks

    def _extract_hyperlink_from_paragraph(self, paragraph) -> str:
        """Extract the first hyperlink URL from a paragraph's runs."""
        try:
            for run in paragraph.runs:
                if hasattr(run, 'hyperlink') and run.hyperlink:
                    address = run.hyperlink.address
                    if address:
                        return address
        except Exception as e:
            logger.debug(f"Could not extract hyperlink: {e}")
        return None

    def _extract_hyperlink_from_text_frame(self, text_frame) -> str:
        """Extract the first hyperlink URL from any paragraph in a text frame."""
        try:
            for p in text_frame.paragraphs:
                url = self._extract_hyperlink_from_paragraph(p)
                if url:
                    return url
        except Exception as e:
            logger.debug(f"Could not extract hyperlink from text frame: {e}")
        return None

    def _extract_video_info(self, shape, slide_order: int, file_id: str, media_dir: str) -> dict:
        """Extract video info from embedded media shape.

        Returns dict with 'url' (external) or 'src' (local saved file) and 'title'.
        """
        try:
            element = None
            if hasattr(shape, '_element'):
                element = shape._element

            if element is None:
                return None

            ns = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main'
            }

            # Look for videoFile in nvPicPr/nvPr (picture shapes)
            videoFile = None
            nvPr = None
            if hasattr(element, 'nvPicPr') and element.nvPicPr is not None:
                nvPr = element.nvPicPr.nvPr
                videoFile = nvPr.find('.//a:videoFile', ns)

            # Also check nvSpPr for other shape types
            if videoFile is None and hasattr(element, 'nvSpPr') and element.nvSpPr is not None:
                nvPr = element.nvSpPr.nvPr
                videoFile = nvPr.find('.//a:videoFile', ns)

            # Check the whole element tree as fallback
            if videoFile is None:
                videoFile = element.find('.//a:videoFile', ns)
                if videoFile is not None:
                    nvPr = element.find('.//p:nvPr', {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})

            if videoFile is None:
                return None

            video_title = shape.name if hasattr(shape, 'name') else "Video"

            # Get the link rId (for external URLs like YouTube)
            video_link_rId = videoFile.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link')
            if video_link_rId:
                target = shape.part.target_ref(video_link_rId)
                # Check if it's an external URL (http/https) or local path
                if target and (target.startswith('http://') or target.startswith('https://')):
                    return {'url': target, 'title': video_title}

            # Check for embedded media (p14:media element for local MP4s)
            if nvPr is not None:
                p14_media = nvPr.find('.//p14:media', ns)
                if p14_media is not None:
                    embed_rId = p14_media.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if embed_rId:
                        # Get the related part (the actual video file)
                        try:
                            video_part = shape.part.related_part(embed_rId)
                            if video_part and hasattr(video_part, 'blob'):
                                # Determine file extension from content type
                                content_type = video_part.content_type if hasattr(video_part, 'content_type') else 'video/mp4'
                                ext_map = {
                                    'video/mp4': '.mp4',
                                    'video/x-m4v': '.m4v',
                                    'video/webm': '.webm',
                                    'video/quicktime': '.mov',
                                    'video/x-msvideo': '.avi',
                                }
                                ext = ext_map.get(content_type, '.mp4')

                                # Save the video file
                                video_filename = f"slide_{slide_order}_{shape.shape_id}{ext}"
                                video_path = os.path.join(media_dir, video_filename)
                                with open(video_path, 'wb') as f:
                                    f.write(video_part.blob)

                                src = f"media/{file_id}/{video_filename}"
                                logger.info(f"Extracted embedded video: {video_filename}")
                                return {'src': src, 'title': video_title}
                        except Exception as e:
                            logger.debug(f"Could not extract embedded video blob: {e}")

        except Exception as e:
            logger.debug(f"Could not extract video info: {e}")
        return None

    def _extract_animation_map(self, slide) -> Dict[int, int]:
        """Extract animation order for shapes on a slide.

        Returns dict mapping shape_id -> animation_order (1-based).
        """
        animation_map: Dict[int, int] = {}
        try:
            # Access the slide's XML element
            slide_elem = slide._element
            ns = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            }

            # Find timing/sequence in slide XML
            timing = slide_elem.find('.//p:timing', ns)
            if timing is None:
                return animation_map

            # Look for sequence of animations
            seq = timing.find('.//p:seq', ns)
            if seq is None:
                return animation_map

            # Each childTnLst contains animation nodes
            order = 0
            for child_tn_lst in seq.findall('.//p:childTnLst', ns):
                for par in child_tn_lst.findall('p:par', ns):
                    # Find the target shape ID
                    for tgt_el in par.findall('.//p:tgtEl', ns):
                        sp_tgt = tgt_el.find('p:spTgt', ns)
                        if sp_tgt is not None:
                            shape_id_str = sp_tgt.get('spid')
                            if shape_id_str:
                                order += 1
                                shape_id = int(shape_id_str)
                                if shape_id not in animation_map:
                                    animation_map[shape_id] = order
        except Exception as e:
            logger.debug(f"Could not extract animation map: {e}")

        return animation_map

    def _extract_formatted_runs(self, paragraph) -> List[TextRun]:
        """Extract ALL text runs with formatting from a paragraph."""
        runs: List[TextRun] = []
        has_any_formatting = False

        try:
            for run in paragraph.runs:
                if not run.text:
                    continue

                # Get formatting
                bold = run.font.bold if run.font.bold is not None else False
                italic = run.font.italic if run.font.italic is not None else False
                underline = run.font.underline is not None and run.font.underline

                # Get URL if present
                url = None
                if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                    url = run.hyperlink.address

                # Get font size (in points)
                font_size = None
                if run.font.size:
                    font_size = run.font.size.pt

                # Get font color
                font_color = None
                try:
                    if run.font.color and run.font.color.rgb:
                        font_color = str(run.font.color.rgb)
                except Exception:
                    pass

                # Track if any run has formatting
                if bold or italic or underline or url or font_color:
                    has_any_formatting = True

                # Include ALL runs to preserve full text
                runs.append(TextRun(
                    text=run.text,
                    bold=bold,
                    italic=italic,
                    underline=underline,
                    url=url,
                    font_size=font_size,
                    font_color=font_color,
                ))

        except Exception as e:
            logger.debug(f"Could not extract formatted runs: {e}")

        # Only return runs if there's any formatting, otherwise return empty
        # (the plain text is already captured in the ListItem.text field)
        return runs if has_any_formatting else []

    def _extract_auto_shape(self, shape, slide, animation_map: Dict[int, int]) -> Optional[ShapeBlock]:
        """Extract auto shape (arrow, connector, symbol, etc.).

        Returns ShapeBlock if this is a meaningful auto shape, None otherwise.
        """
        try:
            # Get shape type
            shape_type_val = shape.shape_type

            # Map shape types we care about
            # MSO_SHAPE_TYPE values: AUTO_SHAPE=1, CALLOUT=2, FREEFORM=5, LINE=9, etc.
            shape_type_names = {
                1: "auto_shape",
                2: "callout",
                5: "freeform",
                9: "line",
                19: "text_box",
                21: "connector",
            }

            # Skip types we handle elsewhere (pictures, tables, groups, etc.)
            skip_types = {6, 7, 13, 14, 16, 17, 18, 19}  # GraphicFrame, Group, Picture, etc.

            if shape_type_val in skip_types:
                return None

            # Get shape name which often describes what it is
            shape_name = shape.name if hasattr(shape, 'name') else ""

            # Check if this looks like a meaningful shape (arrow, symbol, connector)
            meaningful_keywords = [
                'arrow', 'connector', 'line', 'equal', 'plus', 'minus',
                'chevron', 'block', 'star', 'heart', 'lightning', 'sun',
                'callout', 'bubble', 'cloud', 'oval', 'rectangle', 'triangle',
                'pentagon', 'hexagon', 'cross', 'notequal', 'not equal'
            ]

            name_lower = shape_name.lower()
            is_meaningful = any(kw in name_lower for kw in meaningful_keywords)

            # Also check the actual auto shape type if available
            auto_shape_type = None
            if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type:
                auto_shape_type = str(shape.auto_shape_type).split('.')[-1].lower()
                is_meaningful = is_meaningful or auto_shape_type not in ('rectangle', 'rounded_rectangle')

            # Skip plain text boxes and rectangles without special meaning
            if not is_meaningful and shape_type_val == 1:
                # Check if it has visible outline or fill that makes it a shape
                try:
                    if hasattr(shape, 'fill') and shape.fill:
                        fill_type = shape.fill.type
                        if fill_type is None:
                            return None
                except Exception:
                    return None

            if not is_meaningful:
                return None

            # Get position
            left = shape.left if hasattr(shape, 'left') and shape.left else 0
            top = shape.top if hasattr(shape, 'top') and shape.top else 0
            width = shape.width if hasattr(shape, 'width') and shape.width else 0
            height = shape.height if hasattr(shape, 'height') and shape.height else 0

            # Get text inside shape if any
            text = ""
            shape_runs: List[TextRun] = []
            if shape.has_text_frame and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip()
                for p in shape.text_frame.paragraphs:
                    shape_runs.extend(self._extract_formatted_runs(p))

            # Get colors
            fill_color = None
            line_color = None
            try:
                if hasattr(shape, 'fill') and shape.fill and shape.fill.fore_color:
                    if shape.fill.fore_color.rgb:
                        fill_color = str(shape.fill.fore_color.rgb)
            except Exception:
                pass

            try:
                if hasattr(shape, 'line') and shape.line and shape.line.color:
                    if shape.line.color.rgb:
                        line_color = str(shape.line.color.rgb)
            except Exception:
                pass

            # Get rotation
            rotation = 0.0
            if hasattr(shape, 'rotation'):
                rotation = shape.rotation or 0.0

            # Get animation order
            animation_order = None
            shape_id = shape.shape_id if hasattr(shape, 'shape_id') else None
            if shape_id and shape_id in animation_map:
                animation_order = animation_map[shape_id]

            return ShapeBlock(
                shape_type=auto_shape_type or shape_type_names.get(shape_type_val, "shape"),
                shape_name=shape_name,
                text=text,
                runs=shape_runs,
                left=left,
                top=top,
                width=width,
                height=height,
                fill_color=fill_color,
                line_color=line_color,
                rotation=rotation,
                animation_order=animation_order,
            )

        except Exception as e:
            logger.debug(f"Could not extract auto shape: {e}")
            return None
