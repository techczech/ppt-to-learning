import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze():
    f = "sourcefiles/aicc/Introducing ChatGPT and generative AI at Oxford - draft version.pptx"
    if not os.path.exists(f):
        print("File not found!")
        return

    prs = Presentation(f)
    print(f"Analyzing {f}...")
    
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check for MEDIA (16) or generic shape type
            # python-pptx identifies videos often as Picture/GraphicFrame or special Media type depending on version
            
            print(f"Slide {i+1} Shape: {shape.shape_type} ({type(shape).__name__})")
            
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                print("  -> FOUND MEDIA SHAPE!")
                try:
                    # Try accessing format
                    fmt = shape.media_format
                    print(f"     Mime: {fmt.mime_type}")
                    print(f"     Type: {fmt.type}")
                except Exception as e:
                    print(f"     Error accessing media info: {e}")

            # Sometimes videos are Pictures with a link?
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Check if it's a video placeholder
                pass

analyze()
