from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(title_slide_layout)
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Welcome to Interactive Learning"
subtitle.text = "A new way to learn from presentations"

# Slide 2: Content Slide with Bullets
content_slide_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(content_slide_layout)
title = slide2.shapes.title
body = slide2.placeholders[1]

title.text = "Key Concepts"
body.text_frame.text = "Concept 1: Extracted from PowerPoint"
p = body.text_frame.add_paragraph()
p.text = "Concept 2: Rendered as HTML"
p.level = 1
p = body.text_frame.add_paragraph()
p.text = "Concept 3: Enhanced with AI"
p.level = 2

# Slide 3: Content with Speaker Notes
slide3 = prs.slides.add_slide(content_slide_layout)
title = slide3.shapes.title
body = slide3.placeholders[1]

title.text = "Instructor Notes"
body.text_frame.text = "This slide has hidden notes for the instructor."

notes_slide = slide3.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = "This is a speaker note. It provides extra context that isn't on the slide itself, which can be very useful for creating rich learning materials."

# Save the presentation
prs.save("sourcefiles/sample.pptx")

print("Dummy presentation 'sourcefiles/sample.pptx' created successfully.")
